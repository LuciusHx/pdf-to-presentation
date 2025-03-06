import pdfplumber
import re
from pptx import Presentation
from pptx.util import Inches

import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from ttkthemes import ThemedTk 

# ------ selecionar o arquivo PDF ----------
def selecionar_arquivo():
    # janela para o usuário selecionar um arquivo PDF
    caminho_pdf = filedialog.askopenfilename(
        title="Selecione um arquivo PDF",
        filetypes=[("Arquivos PDF", "*.pdf")]
    )
    if caminho_pdf:
        entrada_arquivo.delete(0, tk.END)  # limpa campo de entrada
        entrada_arquivo.insert(0, caminho_pdf)  # insere o caminho do arquivo no campo

# ------ Função para processar o PDF e criar o PowerPoint ----------
def processar_pdf():
    caminho_pdf = entrada_arquivo.get()  # Obtém o caminho do arquivo do campo de entrada
    if not caminho_pdf:
        messagebox.showwarning("Aviso", "Nenhum arquivo PDF selecionado!")
        return

    # cabeça da matriz
    dados_tabela = [["Código", "Especificação", "Valor Empenhado", "Valor Liquidado", "Valor Pago"]]

    def ler_pdf():
        with pdfplumber.open(caminho_pdf) as pdf:
            for page in pdf.pages:
                texto = page.extract_text()
                if texto:
                    linhas = texto.split("\n")
                    for linha in linhas:
                        if re.match(r"^\d{1,2}\.\d{1,2}\.\d{1,2}\.\d{1,2}\.\d{1,2}", linha):
                            match = re.match(r"^(\d{1,2}\.\d{1,2}\.\d{1,2}\.\d{1,2}\.\d{1,2})\s+(.+?)\s+([\d\.,-]+)\s+([\d\.,-]+)\s+([\d\.,-]+)\*?$", linha)
                            if match:
                                codigo = match.group(1)
                                especificacao = match.group(2).strip()
                                empenhado = match.group(3)
                                liquidado = match.group(4)
                                pago = match.group(5)
                                dados_tabela.append([codigo, especificacao, empenhado, liquidado, pago])
                            else:
                                partes = re.split(r"\s{2,}", linha)
                                if len(partes) >= 5:
                                    codigo = partes[0]
                                    especificacao = " ".join(partes[1:-3])
                                    empenhado = partes[-3]
                                    liquidado = partes[-2]
                                    pago = partes[-1].replace("*", "")
                                    dados_tabela.append([codigo, especificacao, empenhado, liquidado, pago])

        if len(dados_tabela) == 1:
            messagebox.showwarning("Aviso", "Nenhum dado válido foi encontrado no PDF.")
            return

    def criar_slides():
        apresentacao = Presentation()
        linhas_por_slide = 15

        for i in range(0, len(dados_tabela), linhas_por_slide):
            slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[5])
            rows = min(linhas_por_slide, len(dados_tabela) - i)
            cols = len(dados_tabela[0])

            left = Inches(0.5)
            top = Inches(0.2)
            width = Inches(9)
            height = Inches(5)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for r, linha in enumerate(dados_tabela[i:i + linhas_por_slide]):
                for c, valor in enumerate(linha):
                    table.cell(r, c).text = valor

        # Salva a apresentação
        pptx_path = "tabela_pptx.pptx"
        apresentacao.save(pptx_path)
        messagebox.showinfo("Sucesso", f"Apresentação criada com sucesso!\nArquivo salvo em: {pptx_path}")

    ler_pdf()
    criar_slides()

# ------ Interface gráfica ----------

# janela principal
janela = ThemedTk(theme="clearlooks")
janela.title("PDF para PPTX")
janela.geometry("500x200")

# estilos
estilo = ttk.Style()
estilo.configure("TButton", font=("Arial", 12), padding=10)
estilo.configure("TEntry", font=("Arial", 12), padding=10)

# organizar os widgets
frame = ttk.Frame(janela)
frame.pack(pady=20, padx=20, fill="both", expand=True)

# Input do arquivo
entrada_arquivo = ttk.Entry(frame, width=40)
entrada_arquivo.grid(row=0, column=0, padx=10, pady=10)

# Botão de selecionar o arquivo
botao_selecionar = ttk.Button(frame, text="Selecionar PDF", command=selecionar_arquivo)
botao_selecionar.grid(row=0, column=1, padx=10, pady=10)

# Botão para processar o PDF
botao_processar = ttk.Button(frame, text="Processar PDF", command=processar_pdf)
botao_processar.grid(row=1, column=0, columnspan=2, pady=10)


#-------------- creditos - rodape --------------
label_desenvolvido = ttk.Label(text="Desenvolvido por Lucius Hebert", font=("Arial", 10))
label_desenvolvido.pack(padx=5, pady=2)

# loop do tkinter
janela.mainloop()